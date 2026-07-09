"""Format-specific text extraction."""

from __future__ import annotations

import csv
import json
import zipfile
from io import BytesIO, StringIO
from pathlib import Path
from typing import Callable

import httpx
from bs4 import BeautifulSoup
from openpyxl import load_workbook
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from app.core.config import get_settings
from app.services.ingestion.cleaning import clean_text
from app.services.ingestion.models import ExtractedDocument, ExtractedPage

Extractor = Callable[[Path], ExtractedDocument]


def _single_page(text: str, *, page_number: int | None = None, **meta) -> ExtractedDocument:
    return ExtractedDocument(
        pages=[ExtractedPage(text=clean_text(text), page_number=page_number)],
        metadata=meta,
    )


def extract_txt(path: Path) -> ExtractedDocument:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return _single_page(text, source_type="text")


def extract_json(path: Path) -> ExtractedDocument:
    raw = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
    text = json.dumps(raw, indent=2, ensure_ascii=False)
    return _single_page(text, source_type="json")


def extract_csv(path: Path) -> ExtractedDocument:
    rows: list[str] = []
    with path.open("r", encoding="utf-8", errors="ignore", newline="") as handle:
        reader = csv.reader(handle)
        for row in reader:
            rows.append(" | ".join(cell.strip() for cell in row))
    return _single_page("\n".join(rows), source_type="csv")


def extract_pdf(path: Path) -> ExtractedDocument:
    reader = PdfReader(str(path))
    pages: list[ExtractedPage] = []
    for index, page in enumerate(reader.pages, start=1):
        text = clean_text(page.extract_text() or "")
        pages.append(ExtractedPage(text=text, page_number=index))

    if not any(p.text for p in pages):
        pages = _ocr_pdf(path).pages

    return ExtractedDocument(pages=pages, metadata={"source_type": "pdf", "ocr_used": not any(p.text for p in pages)})


def _ocr_pdf(path: Path) -> ExtractedDocument:
    settings = get_settings()
    try:
        from pdf2image import convert_from_path
        import pytesseract

        if settings.tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd
        images = convert_from_path(str(path))
        pages = []
        for index, image in enumerate(images, start=1):
            text = clean_text(pytesseract.image_to_string(image))
            pages.append(ExtractedPage(text=text, page_number=index))
        return ExtractedDocument(pages=pages, metadata={"source_type": "pdf", "ocr_used": True})
    except Exception as exc:
        return ExtractedDocument(
            pages=[ExtractedPage(text="", page_number=1)],
            metadata={"source_type": "pdf", "ocr_error": str(exc)},
        )


def extract_docx(path: Path) -> ExtractedDocument:
    document = DocxDocument(str(path))
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    text = "\n".join(paragraphs)
    return _single_page(text, source_type="docx")


def extract_pptx(path: Path) -> ExtractedDocument:
    presentation = Presentation(str(path))
    pages: list[ExtractedPage] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        pages.append(ExtractedPage(text=clean_text("\n".join(parts)), page_number=index))
    return ExtractedDocument(pages=pages, metadata={"source_type": "pptx"})


def extract_xlsx(path: Path) -> ExtractedDocument:
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    pages: list[ExtractedPage] = []
    for index, sheet in enumerate(workbook.worksheets, start=1):
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                rows.append(" | ".join(values))
        pages.append(ExtractedPage(text=clean_text("\n".join(rows)), page_number=index))
    return ExtractedDocument(pages=pages, metadata={"source_type": "xlsx"})


def extract_html(path: Path) -> ExtractedDocument:
    html = path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    return _single_page(soup.get_text("\n"), source_type="html")


def extract_image(path: Path) -> ExtractedDocument:
    settings = get_settings()
    try:
        import pytesseract
        from PIL import Image

        if settings.tesseract_cmd:
            pytesseract.pytesseract.tesseract_cmd = settings.tesseract_cmd
        image = Image.open(path)
        text = clean_text(pytesseract.image_to_string(image))
        return _single_page(text, source_type="image", ocr_used=True)
    except Exception as exc:
        return ExtractedDocument(
            pages=[ExtractedPage(text="", page_number=1)],
            metadata={"source_type": "image", "ocr_error": str(exc)},
        )


def extract_zip(path: Path) -> ExtractedDocument:
    settings = get_settings()
    allowed = settings.allowed_extension_set
    pages: list[ExtractedPage] = []
    metadata: dict = {"source_type": "zip", "members": []}

    with zipfile.ZipFile(path, "r") as archive:
        for member in archive.infolist():
            if member.is_dir():
                continue
            inner_name = Path(member.filename).name
            extension = Path(inner_name).suffix.lower().lstrip(".")
            if extension not in allowed or extension == "zip":
                continue
            data = archive.read(member)
            inner_path = Path(settings.temp_path) / f"zip_{path.stem}_{inner_name}"
            inner_path.parent.mkdir(parents=True, exist_ok=True)
            inner_path.write_bytes(data)
            try:
                inner_doc = extract_file(inner_path, extension=extension)
                for page in inner_doc.pages:
                    pages.append(
                        ExtractedPage(
                            text=f"[{inner_name}]\n{page.text}",
                            page_number=page.page_number,
                        )
                    )
                metadata["members"].append(inner_name)
            finally:
                if inner_path.exists():
                    inner_path.unlink(missing_ok=True)

    if not pages:
        pages = [ExtractedPage(text="", page_number=1)]
    return ExtractedDocument(pages=pages, metadata=metadata)


async def extract_url(url: str) -> ExtractedDocument:
    async with httpx.AsyncClient(timeout=30.0, follow_redirects=True) as client:
        response = await client.get(url)
        response.raise_for_status()
        content_type = response.headers.get("content-type", "")
        if "html" in content_type or url.lower().endswith((".html", ".htm")):
            soup = BeautifulSoup(response.text, "lxml")
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()
            title = (soup.title.string or "").strip() if soup.title else ""
            text = clean_text(soup.get_text("\n"))
            return ExtractedDocument(
                pages=[ExtractedPage(text=text, page_number=1)],
                metadata={"source_type": "url", "url": url, "title": title},
            )
        # Binary / other — persist to temp and extract by extension from URL path
        extension = Path(url.split("?")[0]).suffix.lower().lstrip(".")
        settings = get_settings()
        temp_path = Path(settings.temp_path) / f"url_download{Path(url).suffix or '.bin'}"
        temp_path.parent.mkdir(parents=True, exist_ok=True)
        temp_path.write_bytes(response.content)
        try:
            doc = extract_file(temp_path, extension=extension or None, content_type=content_type)
            doc.metadata["url"] = url
            return doc
        finally:
            temp_path.unlink(missing_ok=True)


_EXTRACTORS: dict[str, Extractor] = {
    "txt": extract_txt,
    "md": extract_txt,
    "markdown": extract_txt,
    "json": extract_json,
    "csv": extract_csv,
    "pdf": extract_pdf,
    "docx": extract_docx,
    "ppt": extract_pptx,
    "pptx": extract_pptx,
    "xls": extract_xlsx,
    "xlsx": extract_xlsx,
    "html": extract_html,
    "htm": extract_html,
    "png": extract_image,
    "jpg": extract_image,
    "jpeg": extract_image,
    "webp": extract_image,
    "tiff": extract_image,
    "zip": extract_zip,
}


def extract_file(
    path: Path,
    *,
    extension: str | None = None,
    content_type: str | None = None,
) -> ExtractedDocument:
    """Dispatch extraction based on file extension."""
    ext = (extension or path.suffix.lower().lstrip(".")).lower()
    if not ext and content_type:
        if "pdf" in content_type:
            ext = "pdf"
        elif "word" in content_type:
            ext = "docx"
        elif "html" in content_type:
            ext = "html"
    extractor = _EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext or 'unknown'}")
    return extractor(path)
